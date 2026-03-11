"""
Intelli-Credit Hackathon PPT Builder
Generates a professional PPTX presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ── Colour Palette ───────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # slide BG
MID_NAVY    = RGBColor(0x1B, 0x26, 0x3B)   # card BG
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xFF)   # primary accent
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)   # secondary accent
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)   # highlight
GREEN       = RGBColor(0x28, 0xA7, 0x45)
AMBER       = RGBColor(0xFF, 0xC1, 0x07)
RED_COL     = RGBColor(0xDC, 0x35, 0x45)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xB0, 0xBE, 0xC5)
TEAL        = RGBColor(0x00, 0xB8, 0xA0)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper utilities ─────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(BLANK)
    fill_slide_bg(s, DARK_NAVY)
    return s

def fill_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1), radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    run.font.name  = "Calibri"
    return txb

def pill(slide, label, l, t, color=ACCENT_BLUE, text_color=WHITE, size=11):
    w = Inches(1.6)
    h = Inches(0.32)
    b = box(slide, l, t, w, h, fill=color)
    txt(slide, label, l, t, w, h, size=size, bold=True,
        color=text_color, align=PP_ALIGN.CENTER)
    return b

def gradient_bar(slide, l, t, w, h, pct, color=GREEN):
    box(slide, l, t, w, h, fill=MID_NAVY, border=ACCENT_BLUE, border_w=Pt(0.5))
    filled = int(w * pct / 100)
    if filled > 0:
        box(slide, l, t, Emu(filled), h, fill=color)
    return

def section_label(slide, text, l=Inches(0.4), t=Inches(0.18)):
    txt(slide, text, l, t, Inches(6), Inches(0.35),
        size=9, color=ACCENT_CYAN, bold=True)

def slide_title(slide, title, subtitle=None):
    txt(slide, title,
        Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.55),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(0.78), Inches(12.3), Inches(0.38),
            size=14, color=LIGHT_GREY)
    # bottom accent line
    box(slide, Inches(0.5), Inches(1.18), Inches(12.33), Inches(0.04),
        fill=ACCENT_BLUE)

def card(slide, l, t, w, h, title=None, title_color=ACCENT_CYAN):
    box(slide, l, t, w, h, fill=MID_NAVY, border=ACCENT_BLUE, border_w=Pt(0.8))
    if title:
        txt(slide, title, l+Inches(0.15), t+Inches(0.1),
            w-Inches(0.3), Inches(0.35), size=11, bold=True, color=title_color)
    return

def bullet_lines(slide, items, l, t, w, size=11, color=WHITE, spacing=0.295):
    for i, item in enumerate(items):
        prefix, rest = ("• ", item) if not item.startswith("  ") else ("  ", item.strip())
        txt(slide, prefix + rest,
            l, Inches(t + i*spacing), w, Inches(0.3),
            size=size, color=color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
s = add_slide()

# top gradient bar
box(s, 0, 0, W, Inches(0.08), fill=ACCENT_BLUE)

# big title block
box(s, Inches(0.5), Inches(1.5), Inches(8.5), Inches(1.3), fill=MID_NAVY, border=ACCENT_BLUE)
txt(s, "INTELLI-CREDIT", Inches(0.7), Inches(1.6), Inches(8), Inches(0.65),
    size=48, bold=True, color=WHITE)
txt(s, "Next-Gen Corporate Credit Appraisal  •  AI-Powered  •  India-First",
    Inches(0.7), Inches(2.22), Inches(8.2), Inches(0.4), size=14, color=ACCENT_CYAN)

# tag line box
box(s, Inches(0.5), Inches(3.0), Inches(8.5), Inches(0.7), fill=ACCENT_BLUE)
txt(s, "Upload docs → Research → Score → CAM Report   in under 90 seconds",
    Inches(0.65), Inches(3.08), Inches(8.2), Inches(0.5),
    size=15, bold=True, color=WHITE)

# 3 pillar badges
cols = [Inches(1.0), Inches(4.0), Inches(7.1)]
labels = ["PILLAR 1\nExtractor", "PILLAR 2\nResearch Agent", "PILLAR 3\nCAM Engine"]
colors = [TEAL, ACCENT_BLUE, ACCENT_GOLD]
t_col = [WHITE, WHITE, DARK_NAVY]
for lx, lab, col, tc in zip(cols, labels, colors, t_col):
    box(s, lx, Inches(4.0), Inches(2.3), Inches(1.0), fill=col)
    txt(s, lab, lx, Inches(4.05), Inches(2.3), Inches(0.95),
        size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)

txt(s, "IIT Hackathon 2026  |  Team Intelli-Credit",
    Inches(0.5), Inches(6.8), Inches(8), Inches(0.38),
    size=11, color=LIGHT_GREY)

# right decorative element
for i in range(6):
    alpha = 0.15 + i*0.13
    r = int(0 + 0*alpha); g = int(122*alpha); b = int(255*alpha)
    box(s, Inches(9.8), Inches(1.2 + i*0.9), Inches(3.2), Inches(0.75),
        fill=RGBColor(min(r,30), min(g,80), min(b,180)))

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "The Problem: India's Credit Data Paradox", "Why corporate credit appraisal is broken today")

# 3 pain point cards
pain = [
    ("⏱  2–3 Weeks",         "Process Time",    "Manual data stitching from\n10+ fragmented sources",           RED_COL),
    ("⚠  Human Bias",         "Decision Quality","No systematic cross-verification\nof declared vs actual data",    AMBER),
    ("📉  Missed Signals",     "Risk Detection",  "India-specific fraud signals\n(GST inflation, ITC fraud) unseen", RED_COL),
]
for i, (stat, label, desc, col) in enumerate(pain):
    lx = Inches(0.4 + i*4.3)
    card(s, lx, Inches(1.4), Inches(3.9), Inches(2.5), title_color=col)
    txt(s, stat,  lx+Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.6),
        size=22, bold=True, color=col)
    txt(s, label, lx+Inches(0.2), Inches(2.12), Inches(3.5), Inches(0.35),
        size=11, bold=True, color=WHITE)
    txt(s, desc,  lx+Inches(0.2), Inches(2.5),  Inches(3.5), Inches(0.7),
        size=10, color=LIGHT_GREY)

# data sources grid
box(s, Inches(0.4), Inches(4.1), Inches(12.53), Inches(0.04), fill=ACCENT_BLUE)
txt(s, "DISPARATE DATA SOURCES A CREDIT MANAGER MUST MANUALLY STITCH",
    Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.3), size=9, color=ACCENT_CYAN, bold=True)

sources = [
    ("Structured", ["GST Filings", "ITR", "Bank Statements"], TEAL),
    ("Unstructured", ["Annual Reports", "Rating Reports", "Board Minutes"], ACCENT_BLUE),
    ("External Intel", ["MCA Filings", "eCourts Portal", "News Reports"], ACCENT_GOLD),
    ("Primary", ["Site Visits", "Mgmt Interviews", "Due Diligence"], GREEN),
]
for i, (cat, items, col) in enumerate(sources):
    lx = Inches(0.4 + i*3.25)
    card(s, lx, Inches(4.55), Inches(3.0), Inches(2.65), title=cat, title_color=col)
    bullet_lines(s, items, lx+Inches(0.2), 5.1, Inches(2.7), size=10, spacing=0.32)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Our Solution: Intelli-Credit AI Engine", "Three AI pillars → One unified Credit Appraisal Memorandum")

pillars = [
    ("01", "DOCUMENT\nEXTRACTOR",   "PILLAR 1", TEAL,
     ["PDF / DOCX / XLSX / CSV", "6 auto-classified doc types",
      "LLM structured extraction", "11 cross-validation checks",
      "CV-009: GST vs Bank fraud ★", "CV-011: GSTR-2A vs 3B ★"]),
    ("02", "RESEARCH\nAGENT",       "PILLAR 2", ACCENT_BLUE,
     ["6 concurrent sources", "RBI Wilful Defaulter list",
      "MCA21 filings + charges", "eCourts litigation check",
      "Tavily news intelligence", "GSTN + CIBIL scores"]),
    ("03", "CAM\nENGINE",           "PILLAR 3", ACCENT_GOLD,
     ["8-dimension scoring model", "22 financial scoring rules",
      "9 cross-pillar contradictions", "LLM narrative generation",
      "Explainable rate derivation", "DOCX + PDF output ★"]),
]
for i, (num, name, label, col, pts) in enumerate(pillars):
    lx = Inches(0.35 + i*4.33)
    # header bar
    box(s, lx, Inches(1.4), Inches(4.0), Inches(0.55), fill=col)
    txt(s, f"{num}  {label}", lx+Inches(0.12), Inches(1.46),
        Inches(3.8), Inches(0.45), size=11, bold=True,
        color=DARK_NAVY if col == ACCENT_GOLD else WHITE)
    # card
    card(s, lx, Inches(1.95), Inches(4.0), Inches(4.5))
    txt(s, name, lx+Inches(0.15), Inches(2.05), Inches(3.7), Inches(0.75),
        size=17, bold=True, color=col)
    bullet_lines(s, pts, lx+Inches(0.15), 2.85, Inches(3.7), size=11, spacing=0.355)

# arrow connectors
for lx in [Inches(4.35), Inches(8.68)]:
    txt(s, "→", lx, Inches(3.7), Inches(0.4), Inches(0.5),
        size=24, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# bottom output strip
box(s, Inches(0.35), Inches(6.65), Inches(12.63), Inches(0.65), fill=MID_NAVY, border=ACCENT_GOLD)
txt(s, "OUTPUT: Downloadable CAM Report (DOCX + PDF)  •  Decision: APPROVE / CONDITIONAL / REJECT  •  Score /100  •  Rate derivation",
    Inches(0.55), Inches(6.75), Inches(12.3), Inches(0.45),
    size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — DIFFERENTIATORS
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "What No Other Team Will Have", "India-specific intelligence baked into the credit score")

diffs = [
    ("★ CV-009", "GST vs Bank Circular Trading Detector",
     "GST declared turnover ÷ Bank credits > 1.30x = CRITICAL fraud flag.\nThe borrower is inflating revenue in GST returns vs actual bank flows.\nIndia-specific  •  baked into the composite score  •  in the CAM doc",
     RED_COL, "CRITICAL FLAG"),
    ("★ CV-011", "GSTR-2A vs GSTR-3B ITC Reconciliation",
     "GSTR-2A is auto-populated from supplier filings — cannot be falsified.\nIf GSTR-3B ITC > GSTR-2A ITC by >10% = Bogus invoice / fabricated credits.\nOnly a team that understands Indian GST architecture builds this.",
     RED_COL, "CRITICAL FLAG"),
    ("★ 9-Rule", "Cross-Pillar Contradiction Detector",
     "Strong financials but adverse research signals? The system explicitly flags\nthis — exactly the scenario in the problem statement example.\nEvery conflict shown in the CAM with source pillar labels.",
     ACCENT_GOLD, "UNIQUE SIGNAL"),
    ("★ ±15pt", "Primary Insight Score Adjustment",
     "Factory visit observations, management quality, CIBIL score feed directly\ninto the composite score as quantitative ±15 point adjustment.\nField observations are NOT just notes — they move the needle.",
     TEAL, "QUANTIFIED"),
]
for i, (badge, title, desc, col, tag) in enumerate(diffs):
    row = i // 2
    col_idx = i % 2
    lx = Inches(0.4 + col_idx*6.5)
    ty = Inches(1.45 + row*2.85)
    card(s, lx, ty, Inches(6.15), Inches(2.55), title_color=col)
    # badge pill
    box(s, lx+Inches(0.15), ty+Inches(0.12), Inches(0.9), Inches(0.38), fill=col)
    txt(s, badge, lx+Inches(0.15), ty+Inches(0.12), Inches(0.9), Inches(0.38),
        size=9, bold=True, color=DARK_NAVY if col==ACCENT_GOLD else WHITE, align=PP_ALIGN.CENTER)
    # tag
    box(s, lx+Inches(5.0), ty+Inches(0.12), Inches(1.0), Inches(0.38), fill=MID_NAVY, border=col)
    txt(s, tag, lx+Inches(5.0), ty+Inches(0.12), Inches(1.0), Inches(0.38),
        size=8, bold=True, color=col, align=PP_ALIGN.CENTER)

    txt(s, title, lx+Inches(1.15), ty+Inches(0.12), Inches(3.8), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txt(s, desc, lx+Inches(0.15), ty+Inches(0.65), Inches(5.8), Inches(1.7),
        size=10, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — PIPELINE FLOW
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "End-to-End Pipeline Architecture", "From document upload to downloadable CAM — under 90 seconds")

# Timeline flow
steps = [
    ("UPLOAD",     "Credit officer\nuploads PDFs,\nDOCX, XLSX",     TEAL,        "~0s"),
    ("EXTRACT",    "Format detect\nLLM structure\n11 validations",   TEAL,        "~15–30s"),
    ("RESEARCH",   "6 concurrent\nasync sources\nRBI+MCA+Courts",    ACCENT_BLUE, "~8–20s"),
    ("SCORE",      "8-dimension\nweighted model\n+contradictions",   ACCENT_BLUE, "~3s"),
    ("NARRATIVE",  "LLM generates\n8 CAM sections\nfull rationale",  ACCENT_GOLD, "~10–20s"),
    ("DOWNLOAD",   "DOCX + PDF\nCAM Report\nready to share",        GREEN,       "~3s"),
]
bw = Inches(1.9)
for i, (title, desc, col, timing) in enumerate(steps):
    lx = Inches(0.3 + i*2.15)
    ty = Inches(1.5)
    # box
    box(s, lx, ty, bw, Inches(2.8), fill=MID_NAVY, border=col, border_w=Pt(2))
    # top color bar
    box(s, lx, ty, bw, Inches(0.45), fill=col)
    txt(s, title, lx, ty, bw, Inches(0.45),
        size=11, bold=True, color=DARK_NAVY if col==ACCENT_GOLD else WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, lx+Inches(0.1), ty+Inches(0.52), bw-Inches(0.2), Inches(1.7),
        size=9.5, color=LIGHT_GREY)
    # timing badge
    box(s, lx+Inches(0.25), ty+Inches(2.35), Inches(1.4), Inches(0.32), fill=col)
    txt(s, timing, lx+Inches(0.25), ty+Inches(2.35), Inches(1.4), Inches(0.32),
        size=9, bold=True, color=DARK_NAVY if col==ACCENT_GOLD else WHITE, align=PP_ALIGN.CENTER)
    # arrow (not last)
    if i < len(steps)-1:
        txt(s, "▶", lx+bw, ty+Inches(1.1), Inches(0.25), Inches(0.5),
            size=13, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# total time banner
box(s, Inches(0.3), Inches(4.55), Inches(12.73), Inches(0.6), fill=ACCENT_BLUE)
txt(s, "TOTAL END-TO-END:  ~45–90 seconds   vs   2–3 weeks manually   =  97% time reduction",
    Inches(0.5), Inches(4.65), Inches(12.5), Inches(0.4),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# tech stack row
techs = ["Next.js 15", "FastAPI", "SQLite → Databricks", "pdfplumber + OCR",
         "httpx asyncio", "Anthropic Claude / Gemini", "python-docx + reportlab"]
txt(s, "TECH STACK:", Inches(0.3), Inches(5.4), Inches(2), Inches(0.3),
    size=9, bold=True, color=ACCENT_CYAN)
for i, t in enumerate(techs):
    box(s, Inches(0.3 + i*1.82), Inches(5.4), Inches(1.75), Inches(0.35),
        fill=MID_NAVY, border=ACCENT_BLUE)
    txt(s, t, Inches(0.32 + i*1.82), Inches(5.42), Inches(1.7), Inches(0.32),
        size=8, color=WHITE, align=PP_ALIGN.CENTER)

# Research agent detail
card(s, Inches(0.3), Inches(5.95), Inches(12.73), Inches(1.35), title="PILLAR 2: Research Agent — 6 Concurrent Sources (asyncio.gather)", title_color=ACCENT_CYAN)
sources6 = ["RBI Wilful\nDefaulter", "MCA21\nFilings", "eCourts\nLitigation", "Tavily\nNews AI", "GSTN\nStatus", "CIBIL\nScores"]
src_col  = [RED_COL, ACCENT_BLUE, AMBER, TEAL, GREEN, ACCENT_BLUE]
for i, (src, col) in enumerate(zip(sources6, src_col)):
    box(s, Inches(0.55 + i*2.1), Inches(6.45), Inches(1.9), Inches(0.65), fill=col)
    txt(s, src, Inches(0.55 + i*2.1), Inches(6.5), Inches(1.9), Inches(0.6),
        size=9, bold=True, color=DARK_NAVY if col in [AMBER] else WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — SCORING MODEL
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "The 8-Dimension Scoring Model", "Composite credit score: deterministic, weighted, explainable")

dimensions = [
    ("Capacity",        25, 78, GREEN,       "DSCR • ICR • CFO • Revenue CAGR"),
    ("Capital",         20, 65, GREEN,       "Net Worth • D/E Ratio • TNW • WC"),
    ("Character",       20, 55, AMBER,       "RBI • MCA • eCourts • News Intel"),
    ("Collateral",      15, 72, GREEN,       "Coverage • Charge rank • Pledge"),
    ("Conditions",      10, 60, AMBER,       "Sector outlook from News AI"),
    ("GST Quality",      5, 81, GREEN,       "Filing compliance • Bank recon"),
    ("Litigation Risk",  3, 70, GREEN,       "eCourts case severity count"),
    ("MCA Compliance",   2, 75, GREEN,       "ROC filing recency • DIN history"),
]
for i, (name, wt, score, col, drivers) in enumerate(dimensions):
    row = i % 4
    c_  = i // 4
    lx  = Inches(0.38 + c_*6.5)
    ty  = Inches(1.42 + row*1.48)
    tw  = Inches(6.15)

    card(s, lx, ty, tw, Inches(1.32))
    txt(s, name,    lx+Inches(0.15), ty+Inches(0.08), Inches(2.8), Inches(0.35), size=12, bold=True, color=WHITE)
    txt(s, f"Weight: {wt}%", lx+Inches(3.2), ty+Inches(0.08), Inches(1.2), Inches(0.35), size=9, color=LIGHT_GREY)
    txt(s, f"{score}/100", lx+Inches(4.5), ty+Inches(0.05), Inches(1.4), Inches(0.4), size=14, bold=True, color=col, align=PP_ALIGN.RIGHT)
    gradient_bar(s, lx+Inches(0.15), ty+Inches(0.52), Inches(5.0), Inches(0.22), score, col)
    txt(s, drivers, lx+Inches(0.15), ty+Inches(0.82), Inches(5.8), Inches(0.3), size=9, color=LIGHT_GREY)

# composite result box
box(s, Inches(0.38), Inches(7.1), Inches(12.57), Inches(0.3),
    fill=MID_NAVY, border=ACCENT_GOLD)
txt(s, "COMPOSITE SCORE: 68/100  →  AMBER  →  CONDITIONAL APPROVAL    |    Interest Rate: 12.75% p.a.    |    Recommended Amount: ₹47.5 Cr",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.28),
    size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — CAM OUTPUT
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "CAM Report Output — IDFC First Bank Format", "Downloadable DOCX + PDF  •  10-section structured memorandum")

sections_ = [
    ("§ 1", "Executive Summary\n& Decision",       ACCENT_GOLD),
    ("§ 2", "Company &\nBusiness Profile",          TEAL),
    ("§ 3", "Character\n(Research Intel)",           RED_COL),
    ("§ 4", "Capacity &\nFinancials",                GREEN),
    ("§ 5", "Capital &\nBalance Sheet",              ACCENT_BLUE),
    ("§ 6", "Collateral\nAnalysis",                  TEAL),
    ("§ 7A","GST Intelligence\n(India-specific ★)",  RED_COL),
    ("§ 8", "Risk Matrix\n& Flags",                  AMBER),
    ("§ 9", "Recommendation\n& Rate Derivation",     GREEN),
    ("§10", "Decision Rationale\n& Audit Trail ★",   ACCENT_GOLD),
]
for i, (num, name, col) in enumerate(sections_):
    row = i // 5
    ci  = i % 5
    lx  = Inches(0.35 + ci*2.6)
    ty  = Inches(1.45 + row*2.6)
    box(s, lx, ty, Inches(2.45), Inches(0.42), fill=col)
    txt(s, num, lx, ty, Inches(2.45), Inches(0.42),
        size=13, bold=True, color=DARK_NAVY if col in [ACCENT_GOLD, AMBER] else WHITE, align=PP_ALIGN.CENTER)
    card(s, lx, ty+Inches(0.42), Inches(2.45), Inches(1.8))
    txt(s, name, lx+Inches(0.1), ty+Inches(0.55), Inches(2.25), Inches(1.5),
        size=10, color=WHITE, align=PP_ALIGN.CENTER)

# right panel — key facts
card(s, Inches(13.05), Inches(1.45), Inches(0.001), Inches(0.001))   # placeholder
# features list
box(s, Inches(2.6*5 + 0.35 + 0.1), Inches(1.45), Inches(2.5), Inches(5.85),    # won't show, skip
    fill=MID_NAVY, border=ACCENT_BLUE)

# bottom strip
box(s, Inches(0.35), Inches(6.75), Inches(12.63), Inches(0.55), fill=MID_NAVY, border=ACCENT_CYAN)
facts = ["IDFC First Bank Grid Style", "python-docx native", "reportlab PDF fallback", "Auto-heals missing files on startup", "~41–45 KB output"]
for i, f in enumerate(facts):
    txt(s, f"✓  {f}", Inches(0.6 + i*2.5), Inches(6.83), Inches(2.4), Inches(0.35),
        size=9.5, bold=True, color=ACCENT_CYAN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — DEMO WALKTHROUGH
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Live Demo — Bhushan Power & Steel Case", "Pre-seeded case CASE_2026_DEMO01 ready as fallback")

steps_demo = [
    ("1", "LOGIN",            "credit_manager@bank.com\npassword: manager123",                         TEAL),
    ("2", "DASHBOARD",        "See case list with\nrisk bands + decisions",                            ACCENT_BLUE),
    ("3", "OPEN BPSL CASE",   "CASE_2026_DEMO01\nBhushan Power & Steel Ltd",                          ACCENT_BLUE),
    ("4", "VIEW CAM",         "Score: 68/100 AMBER\nConditional Approval\n₹47.5 Cr @ 12.75%",         ACCENT_GOLD),
    ("5", "CONTRADICTIONS",   "Show cross-pillar\ncontradiction callouts\n(the differentiator!)",      RED_COL),
    ("6", "DOWNLOAD",         "Click PDF / DOCX\nFull 10-section CAM\nauto-generated",                 GREEN),
]
for i, (num, title, desc, col) in enumerate(steps_demo):
    lx = Inches(0.3 + (i%3)*4.35)
    ty = Inches(1.42 + (i//3)*2.9)
    box(s, lx, ty, Inches(4.1), Inches(0.5), fill=col)
    txt(s, f"STEP {num}  —  {title}", lx+Inches(0.15), ty+Inches(0.08),
        Inches(3.8), Inches(0.38), size=12, bold=True,
        color=DARK_NAVY if col==ACCENT_GOLD else WHITE)
    card(s, lx, ty+Inches(0.5), Inches(4.1), Inches(2.1))
    txt(s, desc, lx+Inches(0.2), ty+Inches(0.65), Inches(3.7), Inches(1.8),
        size=11, color=WHITE)

# fallback note
box(s, Inches(0.3), Inches(7.05), Inches(12.73), Inches(0.32), fill=MID_NAVY, border=AMBER)
txt(s, "⚡ LIVE PIPELINE:  Upload BPSL DOCX → Pipeline runs in ~60s → Download fresh CAM    |    FALLBACK: CASE_2026_DEMO01 always ready",
    Inches(0.5), Inches(7.08), Inches(12.5), Inches(0.28),
    size=9.5, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — NUMBERS THAT MATTER
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Measurable Claims — Real Numbers Only", "Every number traceable to the actual source code")

stats = [
    ("< 90s",    "End-to-end pipeline\nUpload → CAM Report",    ACCENT_GOLD, "2–3 weeks manually"),
    ("6",        "Concurrent research\nsources via asyncio",     ACCENT_BLUE, "All run in parallel"),
    ("11",       "Cross-validation\nchecks (CV_001–011)",        TEAL,        "11 deterministic rules"),
    ("22",       "Financial scoring\nrules (ScoreBreakdown)",    GREEN,       "Capacity+Capital+Collateral+GST"),
    ("9",        "Cross-pillar\ncontradiction rules",            RED_COL,     "Natural language output"),
    ("4",        "Input formats\nPDF/DOCX/XLSX/CSV",            ACCENT_CYAN, "6 auto-classified doc types"),
    ("±15 pts",  "Primary Insight\nadjustment range",           AMBER,       "Field obs → score impact"),
    ("~12.5K",   "Lines of code\nacross 50 source files",       ACCENT_BLUE, "Production-grade codebase"),
]
for i, (val, label, col, sub) in enumerate(stats):
    row = i // 4
    ci  = i % 4
    lx  = Inches(0.35 + ci*3.25)
    ty  = Inches(1.45 + row*2.85)
    card(s, lx, ty, Inches(3.05), Inches(2.55))
    txt(s, val,   lx+Inches(0.15), ty+Inches(0.12), Inches(2.75), Inches(0.85), size=30, bold=True, color=col)
    txt(s, label, lx+Inches(0.15), ty+Inches(0.98), Inches(2.75), Inches(0.75), size=11, color=WHITE)
    box(s, lx+Inches(0.15), ty+Inches(1.88), Inches(2.75), Inches(0.38), fill=DARK_NAVY)
    txt(s, sub,   lx+Inches(0.2),  ty+Inches(1.9),  Inches(2.65), Inches(0.35), size=8, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — PRODUCTION READINESS
# ═══════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Production Readiness & Databricks Story", "One environment variable separates demo from enterprise scale")

# left: feature checklist
card(s, Inches(0.35), Inches(1.42), Inches(6.0), Inches(5.9), title="✅  What's Production-Ready", title_color=GREEN)
checks = [
    "JWT authentication (24h tokens, role-based)",
    "SQLite → Databricks Delta Lake (USE_DATABRICKS=true)",
    "Startup auto-healer: regenerates missing DOCX/PDF on boot",
    "On-demand document regeneration: POST /cam/regenerate",
    "All paths stored as absolute (no relative path bugs)",
    "3-strategy PDF: docx2pdf → LibreOffice → reportlab",
    "6 research sources with retry + exponential backoff",
    "Graceful degradation: DOCX served if PDF conversion fails",
    "Cross-pillar contradiction audit trail in every CAM",
    "Full explainability chain: every basis point justified",
]
bullet_lines(s, checks, Inches(0.55), 1.95, Inches(5.65), size=10.5, spacing=0.355)

# right: databricks diagram
card(s, Inches(6.6), Inches(1.42), Inches(6.38), Inches(5.9), title="🔄  Databricks Integration Path", title_color=ACCENT_BLUE)

envs = [
    ("HACKATHON MODE",  "USE_DATABRICKS=false", "SQLite • Local files • JSON fallbacks",     TEAL),
    ("PRODUCTION MODE", "USE_DATABRICKS=true",  "Databricks Delta Lake • DBFS • MLflow",     ACCENT_GOLD),
]
for i, (mode, env, desc, col) in enumerate(envs):
    ty_ = 2.1 + i*1.5
    box(s, Inches(6.85), Inches(ty_), Inches(5.85), Inches(1.2), fill=DARK_NAVY, border=col, border_w=Pt(2))
    txt(s, mode, Inches(7.0), Inches(ty_+0.08), Inches(4), Inches(0.35), size=12, bold=True, color=col)
    txt(s, env,  Inches(7.0), Inches(ty_+0.45), Inches(4), Inches(0.3), size=10, color=WHITE, italic=True)
    txt(s, desc, Inches(7.0), Inches(ty_+0.78), Inches(4), Inches(0.3), size=9, color=LIGHT_GREY)

txt(s, "↕  ONE ENV VAR", Inches(8.5), Inches(3.6), Inches(3), Inches(0.4),
    size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# databricks capabilities
cap = ["Databricks Model Serving (LLM inference)", "Delta Lake (case + audit storage)",
       "MLflow experiment tracking", "Unity Catalog (data governance)"]
txt(s, "Full Databricks Stack:", Inches(6.85), Inches(5.05), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=ACCENT_CYAN)
bullet_lines(s, cap, Inches(6.85), 5.35, Inches(5.65), size=10, spacing=0.33)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_slide_bg(s, DARK_NAVY)

box(s, 0, 0, W, Inches(0.08), fill=ACCENT_GOLD)
box(s, 0, H-Inches(0.08), W, Inches(0.08), fill=ACCENT_GOLD)

txt(s, "INTELLI-CREDIT", Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.2),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "The only credit appraisal system that cross-verifies GST declared revenue\nagainst bank credits to detect circular trading — India's #1 lending fraud signal.",
    Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.1),
    size=15, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

box(s, Inches(2.5), Inches(3.8), Inches(8.33), Inches(0.06), fill=ACCENT_GOLD)

bullets_closing = [
    "★  GST vs Bank fraud detector  (CV-009)  — baked into the credit score",
    "★  GSTR-2A vs GSTR-3B ITC reconciliation  (CV-011)  — fabricated input tax credit",
    "★  9-rule cross-pillar contradiction detector  — explicit conflict disclosure",
    "★  Full explainability chain  — every basis point in the rate is justified",
    "★  Under 90 seconds  — from upload to downloadable CAM report (PDF + DOCX)",
]
for i, b in enumerate(bullets_closing):
    txt(s, b, Inches(2.0), Inches(4.1 + i*0.52), Inches(9.3), Inches(0.45),
        size=12, color=WHITE if i>0 else ACCENT_GOLD, bold=(i==0))

txt(s, "IIT Hackathon 2026  •  Team Intelli-Credit",
    Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.38),
    size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────
out = Path(r"C:\Users\venka_5gwzxwk\OneDrive\Desktop\IIT-Hackathon\Intelli_Credit_Presentation.pptx")
prs.save(str(out))
print(f"[OK] Saved: {out}")
print(f"     Slides: {len(prs.slides)}")
print(f"     Size: {out.stat().st_size:,} bytes")
